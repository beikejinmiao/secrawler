#!/usr/bin/env python
# -*- coding:utf-8 -*-
import os
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
import win32com.client


# import textract
# text = textract.process(r'E:\tmp\zdump\test.pptx')
# print(text.decode('utf-8'))       # 无法获取表格内容


def docx(filepath):
    texts = list()
    fopen = Document(filepath)
    for p in fopen.paragraphs:
        texts.extend(p.text)
    for table in fopen.tables:
        for row in table.rows:
            for cell in row.cells:
                texts.extend(cell.text)
    return set(texts)


def pptx(filepath):
    texts = list()
    fopen = Presentation(filepath)
    for slide in fopen.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.extend(shape.text)
            if shape.has_table:
                tbl = shape.table
                for r in range(0, len(tbl.rows)):
                    for c in range(0, len(tbl.columns)):
                        cell = tbl.cell(r, c)
                        texts.extend(cell.text_frame.text)
    return set(texts)


def xlsx(filepath):
    texts = list()
    fopen = pd.read_excel(filepath, sheet_name=None)
    for name, sheet in fopen.items():
        for index, row in sheet.iterrows():
            for value in row:
                texts.extend(str(value))
    return set(texts)


def doc(filepath):
    # https://www.cnblogs.com/zhuminghui/p/11765401.html
    texts = list()
    app = win32com.client.Dispatch('Word.Application')
    # 如果不声明下列属性，运行的时候会显示的打开office软件操作文档
    app.Visible = 0            # 后台运行
    app.DisplayAlerts = 0      # 不显示，不警告
    fopen = app.Documents.Open(filepath)
    for paragraph in fopen.Paragraphs:
        texts.extend(paragraph.Range.Text)
    for table in fopen.Tables:
        for row in table.Rows:
            for cell in row.Cells:
                texts.extend(cell.Range.Text)
    fopen.Close()
    app.Quit()
    return set(texts)


def ppt(filepath):
    app = win32com.client.DispatchEx('Powerpoint.Application')
    ppt.DisplayAlerts = 0                   # 不显示,不警告
    fopen = app.Presentations.Open(filepath)
    new_pptx_file = filepath + '.pptx'
    fopen.SaveAs(new_pptx_file, 24)          # 保存为pptx
    # fopen.SaveAs(filepath + '.pdf', 32)    # 保存为pdf
    fopen.Close()
    app.Quit()
    texts = pptx(new_pptx_file)
    os.remove(new_pptx_file)
    return texts


# def ppt(filepath):
#     ppt = win32com.client.DispatchEx('Powerpoint.Application')
#     ppt.DisplayAlerts = 0  # 不显示,不警告
#     file = ppt.Presentations.Open(filepath)
#     # 用pywin32实在不知道怎么访问ppt文本和表格内容
#     for slide in file.Slides:
#         for shape in slide.Shapes:
#             # https://docs.microsoft.com/zh-cn/office/vba/api/powerpoint.shape.textframe
#             # https://docs.microsoft.com/zh-cn/office/vba/api/powerpoint.textrange
#             # https://docs.microsoft.com/zh-cn/office/vba/api/powerpoint.textrange.text
#             print(shape.TextFrame.TextRange)                        # 任意文本框里的内容
#             # Exception: pywintypes.com_error: (-2147352567, '发生意外。', (0, None, None, None, 0, -2147467259), None)
#             # import win32api
#             # win32api.FormatMessage(-2147467259)
#             # >> '未指定的错误\r\n'
#             # https://blog.csdn.net/huicheng0703/article/details/122242972
#             # https://www.cnblogs.com/vhills/p/8098715.html
#             # for paragraph in shape.TextFrame.TextRange.Paragraphs:
#             #     print(paragraph.Text)


def xls(filepath):
    return xlsx(filepath)


def pdf(filepath):
    texts = list()
    with pdfplumber.open(filepath) as fopen:
        for page in fopen.pages:
            text = page.extract_text()  # 提取文本
            texts.extend(text)
    return set(texts)


