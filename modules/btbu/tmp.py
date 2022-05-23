#!/usr/bin/env python
# -*- coding:utf-8 -*-
import re
import requests
import os
import json
import wget
import pandas as pd
from collections import Counter
from docx import Document
import rarfile
import pdfplumber
from id_validator import validator
from urllib.parse import urlparse
from modules.btbu.util import find_idcard
from utils import reader
from paths import DUMP_HOME


# with open('test.pptx', 'rb') as fopen:
#     for line in fopen.readlines():
#         texts = re.findall(rb'[\x20-\x7e]+', line)
#         print('>>', line)
#         for t in texts:
#             print(str(t, encoding='utf-8'))


# doc = Document('test.docx')
# for p in doc.paragraphs:
#     print('>>', p.text)
#

# default_headers = {
#     'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/99.0.4844.74 Safari/537.36',
#     'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
# }
# resp = requests.get('http://www.btbu.edu.cn/eapdomain/static/cmsfiles/pim/File/20090914095604.doc', allow_redirects=True, headers=default_headers)
# with open('xxx.doc', 'wb') as fopen:
#     fopen.write(resp.content)
# print(resp.status_code)
# print(resp.url)
# print(resp.history)
# print(str(resp.content, encoding='utf-8'))

#
# url = 'http://www.btbu.edu.cn/eapdomain/static/cmsfiles/pim/File/20090914095604.doc'
# filename = wget.download(url)


# xls = pd.read_excel(r'D:\PycharmProjects\secrawler\zdump\downloads\20100610014824.xls', sheet_name=None)
# # xls = pd.read_excel(os.path.join(DUMP_HOME, 'btbu_results.20220426.xlsx'), sheet_name=None)
# for name, sheet in xls.items():
#     # for index, row in sheet.iterrows():
#     #     for value in row:
#     #         print(value)
#     if name == 'all_urls':
#         sheet['is_valid'] = sheet['idcard'].map(lambda x: 1 if validator.is_valid(x) else 0)
#         sheet.to_csv(os.path.join(DUMP_HOME, 'btbu_results.20220426.csv'), index=False)

#
# suffix = list()
# for url in reader(os.path.join(DUMP_HOME, 'file_urls.20220426.txt')):
#     suffix.append(urlparse(url).path.split('.')[-1].lower())
#
# counter = Counter(suffix)
# print(json.dumps(dict(counter.most_common()), indent=4))

# rarfile.UNRAR_TOOL = r"C:\OptSoft\unrar\UnRAR.exe"
# rar = rarfile.RarFile('D:\\PycharmProjects\\secrawler\\zdump\\downloads\\20100611024650.rar')
# with rar as rf:
#     rf.extractall('D:\\var')


# with pdfplumber.open(r'D:\var\蒋茜交强.pdf') as pdf:
#     for page in pdf.pages:
#         text = page.extract_text()  # 提取文本
#         print(find_idcard(text))

from pptx import Presentation

prs = Presentation(r'D:\var\test.pptx')
for slide in prs.slides:
    for shape in slide.shapes:
        print(shape.text)


